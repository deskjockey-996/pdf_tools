from flask import Blueprint, request, send_file, render_template
import os
import tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import zipfile
from werkzeug.utils import secure_filename

pdf2ppt = Blueprint('pdf2ppt', __name__)

@pdf2ppt.route('/')
def pdf2ppt_page():
    return render_template('pdf2ppt.html')

@pdf2ppt.route('', methods=['POST'])
def convert_pdf_to_ppt():
    if 'pdfFiles' not in request.files:
        return {'error': '未找到文件'}, 400
    
    files = request.files.getlist('pdfFiles')
    if not files or files[0].filename == '':
        return {'error': '未选择文件'}, 400

    try:
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            # 创建ZIP文件
            zip_path = os.path.join(temp_dir, 'converted_ppts.zip')
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for file in files:
                    if not file.filename.lower().endswith('.pdf'):
                        continue

                    # 保存上传的PDF
                    pdf_path = os.path.join(temp_dir, secure_filename(file.filename))
                    file.save(pdf_path)

                    # 创建新的PPT演示文稿
                    prs = Presentation()
                    
                    # 设置幻灯片尺寸为16:9
                    prs.slide_width = Inches(16)
                    prs.slide_height = Inches(9)

                    # 打开PDF文件
                    pdf_document = fitz.open(pdf_path)
                    
                    for page_num in range(len(pdf_document)):
                        # 获取PDF页面
                        page = pdf_document[page_num]
                        
                        # 创建新的幻灯片
                        slide_layout = prs.slide_layouts[6]  # 空白布局
                        slide = prs.slides.add_slide(slide_layout)

                        # 将PDF页面渲染为高质量图片
                        zoom = 2  # 提高分辨率
                        mat = fitz.Matrix(zoom, zoom)
                        pix = page.get_pixmap(matrix=mat)
                        
                        # 保存为临时图片文件
                        img_path = os.path.join(temp_dir, f'page_{page_num + 1}.png')
                        pix.save(img_path)
                        
                        # 添加图片到幻灯片（全屏显示）
                        slide.shapes.add_picture(
                            img_path,
                            0, 0,
                            width=prs.slide_width,
                            height=prs.slide_height
                        )

                        # 删除临时图片文件
                        os.remove(img_path)

                    # 保存PPT文件
                    pptx_filename = secure_filename(file.filename.replace('.pdf', '.pptx'))
                    pptx_path = os.path.join(temp_dir, pptx_filename)
                    prs.save(pptx_path)
                    
                    # 添加到ZIP文件
                    zipf.write(pptx_path, pptx_filename)
                    
                    # 关闭PDF文件
                    pdf_document.close()
                    
                    # 删除临时文件
                    os.remove(pptx_path)
                    os.remove(pdf_path)

            # 返回ZIP文件
            return send_file(
                zip_path,
                as_attachment=True,
                download_name='converted_ppts.zip',
                mimetype='application/zip'
            )

    except Exception as e:
        return {'error': f'转换失败: {str(e)}'}, 500 